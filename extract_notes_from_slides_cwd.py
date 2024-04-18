import os
from pptx import Presentation

def extract_speaker_notes(pptx_file, output_file):
    # Load the PPTX file
    prs = Presentation(pptx_file)

    # Open the output file in write mode
    with open(output_file, 'w', encoding='utf-8') as f:
        # Iterate over each slide in the presentation
        for slide in prs.slides:
            # Write the slide number to the output file
            f.write(f"Slide {slide.slide_id}\n")

            # Get the speaker notes for the slide
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide else ""

            # Write the speaker notes to the output file
            f.write(notes_text.strip() + "\n")

            # Add a separator between slides
            f.write("\n---\n\n")

    print(f"Speaker notes extracted successfully to {output_file}")

# Get the current working directory
current_directory = os.getcwd()

print(current_directory)

# Iterate over all files in the current working directory
for filename in os.listdir(current_directory):
    # Check if the file has a .pptx extension
    if filename.endswith('.pptx'):
        # Construct the full path to the input PPTX file
        pptx_file_path = os.path.join(current_directory, filename)

        # Construct the output file path with [Notes] at the beginning
        output_file_path = os.path.join(current_directory, f"[Notes] {os.path.splitext(filename)[0]}.txt")

        # Call the function to extract speaker notes
        extract_speaker_notes(pptx_file_path, output_file_path)